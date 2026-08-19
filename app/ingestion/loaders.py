"""
Extração de texto por formato de arquivo (Etapa 2 do pipeline: Processamento e Extração).

Cada função recebe o caminho de um arquivo e retorna uma string de texto limpo,
pronta para ser dividida em chunks na etapa seguinte.
"""
import json
import re
from pathlib import Path

import pandas as pd
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader


def _limpar_texto(texto: str) -> str:
    """Remove espaços duplicados, linhas em branco excessivas e caracteres de controle."""
    texto = re.sub(r"[ \t]+", " ", texto)
    texto = re.sub(r"\n{3,}", "\n\n", texto)
    return texto.strip()


def extrair_pdf(caminho: Path) -> str:
    """Extrai texto de PDFs nativos (gerados digitalmente).

    Observação: para PDFs escaneados (imagem), seria necessário OCR
    (ex: pytesseract + pdf2image). Não incluído aqui para manter o
    escopo do MVP; ver README para instruções de extensão.
    """
    reader = PdfReader(str(caminho))
    partes = []
    for i, page in enumerate(reader.pages):
        texto_pagina = page.extract_text() or ""
        if texto_pagina.strip():
            partes.append(f"[página {i + 1}]\n{texto_pagina}")
    return _limpar_texto("\n\n".join(partes))


def extrair_docx(caminho: Path) -> str:
    """Extrai texto corrido de documentos Word, preservando títulos como marcações."""
    doc = DocxDocument(str(caminho))
    partes = []
    for para in doc.paragraphs:
        if not para.text.strip():
            continue
        if para.style.name.startswith("Heading"):
            partes.append(f"\n## {para.text}\n")
        else:
            partes.append(para.text)
    # Tabelas também podem conter conteúdo relevante
    for tabela in doc.tables:
        for linha in tabela.rows:
            celulas = [c.text.strip() for c in linha.cells]
            if any(celulas):
                partes.append(" | ".join(celulas))
    return _limpar_texto("\n".join(partes))


def extrair_xlsx(caminho: Path) -> str:
    """Converte planilhas em texto estruturado, linha a linha, repetindo o cabeçalho."""
    xls = pd.ExcelFile(caminho)
    partes = []
    for aba in xls.sheet_names:
        df = xls.parse(aba).fillna("")
        partes.append(f"[Planilha: {aba}]")
        colunas = list(df.columns)
        for _, linha in df.iterrows():
            descricao = "; ".join(f"{col}: {linha[col]}" for col in colunas)
            partes.append(descricao)
    return _limpar_texto("\n".join(partes))


def extrair_pptx(caminho: Path) -> str:
    """Extrai texto de cada slide, incluindo as notas do apresentador."""
    prs = Presentation(str(caminho))
    partes = []
    for i, slide in enumerate(prs.slides):
        textos_slide = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                textos_slide.append(shape.text_frame.text)
        notas = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notas = slide.notes_slide.notes_text_frame.text.strip()
        bloco = f"[Slide {i + 1}]\n" + "\n".join(textos_slide)
        if notas:
            bloco += f"\n(Notas do apresentador: {notas})"
        partes.append(bloco)
    return _limpar_texto("\n\n".join(partes))


def extrair_markdown(caminho: Path) -> str:
    """Markdown já é texto legível; apenas remove marcações técnicas mais pesadas."""
    texto = caminho.read_text(encoding="utf-8")
    texto = re.sub(r"```.*?```", "", texto, flags=re.DOTALL)  # blocos de código
    texto = re.sub(r"[#>*_`]", "", texto)  # símbolos de formatação
    return _limpar_texto(texto)


def extrair_csv(caminho: Path) -> str:
    """Converte CSV em frases legíveis, uma por linha, com cabeçalho repetido."""
    df = pd.read_csv(caminho).fillna("")
    partes = []
    colunas = list(df.columns)
    for _, linha in df.iterrows():
        descricao = "; ".join(f"{col}: {linha[col]}" for col in colunas)
        partes.append(descricao)
    return _limpar_texto("\n".join(partes))


def extrair_json(caminho: Path) -> str:
    """Achata um JSON em frases descritivas do tipo 'chave: valor'."""
    dados = json.loads(caminho.read_text(encoding="utf-8"))

    def achatar(obj, prefixo=""):
        linhas = []
        if isinstance(obj, dict):
            for chave, valor in obj.items():
                novo_prefixo = f"{prefixo}.{chave}" if prefixo else chave
                linhas.extend(achatar(valor, novo_prefixo))
        elif isinstance(obj, list):
            for i, item in enumerate(obj):
                linhas.extend(achatar(item, f"{prefixo}[{i}]"))
        else:
            linhas.append(f"{prefixo}: {obj}")
        return linhas

    return _limpar_texto("\n".join(achatar(dados)))


def extrair_html(caminho: Path) -> str:
    """Remove tags HTML e extrai apenas o texto visível."""
    soup = BeautifulSoup(caminho.read_text(encoding="utf-8"), "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return _limpar_texto(soup.get_text(separator="\n"))


EXTRATORES = {
    ".pdf": extrair_pdf,
    ".docx": extrair_docx,
    ".xlsx": extrair_xlsx,
    ".pptx": extrair_pptx,
    ".md": extrair_markdown,
    ".csv": extrair_csv,
    ".json": extrair_json,
    ".html": extrair_html,
    ".htm": extrair_html,
}


def extrair_texto(caminho: Path) -> str:
    """Ponto de entrada único: escolhe o extrator certo pela extensão do arquivo."""
    extensao = caminho.suffix.lower()
    if extensao not in EXTRATORES:
        raise ValueError(f"Formato não suportado: {extensao}")
    return EXTRATORES[extensao](caminho)
